# ==============================================================================
# COMPLETE CYBERSECURITY AI OFFER GENERATOR PROJECT
# Enhanced version with full pipeline logic from notebook
# ==============================================================================

# app.py - Main Flask Application
import os
import json
import time
from flask import Flask, render_template, request, jsonify, Response
from flask_cors import CORS
import asyncio
from threading import Thread
import traceback
import psutil  # Added for memory usage tracking
import re  # Added for quality metrics calculation
from sklearn.metrics.pairwise import cosine_similarity

# Import your existing modules (with fixed imports)
from langsmith import Client, traceable
from langsmith.wrappers import wrap_openai
from langchain_core.tracers import LangChainTracer
from langchain.callbacks import StdOutCallbackHandler
from langchain_core.callbacks import CallbackManager
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
from langchain_community.vectorstores.faiss import FAISS
from langchain.docstore.document import Document
from langchain_community.docstore.in_memory import InMemoryDocstore
import numpy as np
import faiss
from langchain_core.messages import HumanMessage
from langchain_core.output_parsers import StrOutputParser
from langchain_cohere import ChatCohere
from langchain_core.prompts import ChatPromptTemplate
from pydantic import BaseModel, Field
from typing import List
from typing_extensions import TypedDict
from langgraph.graph import END, StateGraph, START

# Initialize Flask app
app = Flask(__name__)
CORS(app)
    
# ==============================================================================
# CONFIGURATION
# ==============================================================================

# LangSmith and API Configuration
os.environ["LANGCHAIN_TRACING_V2"] = "true"
os.environ["LANGCHAIN_ENDPOINT"] = "https://api.smith.langchain.com"
os.environ["LANGCHAIN_API_KEY"] = "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"
os.environ["LANGCHAIN_PROJECT"] = "AI_Offer"
os.environ["COHERE_API_KEY"] = "PaMzzuLMXnPWgeLwhFi0O5QLQrGrZ9v84TFB44xA"

# Initialize LangSmith client
try:
    langsmith_client = Client()
    langchain_tracer = LangChainTracer()
    callback_manager = CallbackManager([StdOutCallbackHandler(), langchain_tracer])
    print("✅ LangSmith configuration and tracing setup completed!")
except Exception as e:
    print(f"⚠️ LangSmith setup warning: {e}")
    callback_manager = CallbackManager([StdOutCallbackHandler()])

# ==============================================================================
# TAXONOMY
# ==============================================================================

OFFER_TAXONOMY = {
    "cybersecurity": {
        "threat_detection": ["SIEM", "SOC", "threat_hunting", "anomaly_detection"],
        "vulnerability_management": ["penetration_testing", "security_audit", "risk_assessment"],
        "compliance": ["GDPR", "ISO27001", "SOX", "PCI_DSS"],
        "incident_response": ["forensics", "recovery", "investigation", "containment"],
        "identity_management": ["SSO", "MFA", "privileged_access", "identity_governance"]
    },
    "ai_solutions": {
        "machine_learning": ["supervised", "unsupervised", "reinforcement", "deep_learning"],
        "nlp": ["chatbots", "document_processing", "sentiment_analysis", "language_models"],
        "computer_vision": ["object_detection", "image_classification", "facial_recognition"],
        "predictive_analytics": ["forecasting", "anomaly_detection", "pattern_recognition"]
    }
}

@traceable(name="taxonomy_classification")
def classify_offer_taxonomy(question):
    """Classify offer taxonomy with LangSmith tracing"""
    try:
        question_lower = question.lower()
        classification = {"primary": None, "secondary": [], "confidence": 0}
        
        for category, subcategories in OFFER_TAXONOMY.items():
            if any(subcat in question_lower for subcat_list in subcategories.values() for subcat in subcat_list):
                classification["primary"] = category
                classification["confidence"] = 0.8
                
                for subcat_name, subcat_list in subcategories.items():
                    if any(subcat in question_lower for subcat in subcat_list):
                        classification["secondary"].append(subcat_name)
        
        return classification
    except Exception as e:
        print(f"Error in taxonomy classification: {e}")
        return {"primary": "unknown", "secondary": [], "confidence": 0}

# ==============================================================================
# DOCUMENT PROCESSING
# ==============================================================================

# Load sentence-transformer model
try:
    embedding_model = SentenceTransformer("BAAI/bge-base-en-v1.5")
    print("✅ Embedding model loaded successfully")
except Exception as e:
    print(f"❌ Error loading embedding model: {e}")
    embedding_model = None

# Define PPTX path (UPDATE THIS PATH TO YOUR ACTUAL PATH)
drive_path = r"C:\Users\USER\Downloads\générateur d'offres\interface"

@traceable(name="pptx_text_extraction")
def extract_text_from_pptx(file_path):
    """Extract text from PPTX files with LangSmith tracing"""
    try:
        prs = Presentation(file_path)
        extracted_text = "\n".join(
            shape.text.strip()
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
        return extracted_text
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return ""

@traceable(name="document_loading")
def load_documents(drive_path):
    """Load documents with LangSmith tracing"""
    documents = []
    processed_files = 0
    error_files = 0
    
    if not os.path.exists(drive_path):
        print(f"Warning: Path {drive_path} does not exist. Using empty document set.")
        return documents
    
    for root, _, files in os.walk(drive_path):
        for fname in files:
            if fname.lower().endswith(".pptx") and not fname.startswith("~$"):
                fpath = os.path.join(root, fname)
                try:
                    text = extract_text_from_pptx(fpath)
                    if text.strip():  # Only add if text is not empty
                        documents.append(Document(page_content=text, metadata={"source": fpath}))
                        processed_files += 1
                except Exception as e:
                    print(f"Error reading {fpath}: {e}")
                    error_files += 1
    
    print(f"Processed {processed_files} files, {error_files} errors")
    return documents

@traceable(name="document_chunking")
def split_documents(documents):
    """Split documents into chunks with LangSmith tracing"""
    if not documents:
        print("No documents to split, using empty chunks")
        return []
        
    splitter = RecursiveCharacterTextSplitter(chunk_size=512, chunk_overlap=50)
    doc_chunks = splitter.split_documents(documents)
    
    print(f"Split {len(documents)} documents into {len(doc_chunks)} chunks")
    return doc_chunks

# Initialize document processing
print("🔄 Loading and processing documents...")
documents = load_documents(drive_path)
doc_chunks = split_documents(documents)

# Embedding function
class EmbeddingFunction:
    def __init__(self, model):
        self.model = model
    
    @traceable(name="text_embedding")
    def __call__(self, text):
        if isinstance(text, str):
            return self.embed_query(text)
        else:
            return self.embed_documents(text)
    
    @traceable(name="document_embedding")
    def embed_documents(self, texts):
        if self.model is None:
            return [[0.0] * 768 for _ in texts]  # Fallback embeddings
        embeddings_tensor = self.model.encode(texts, convert_to_tensor=True)
        embeddings_np = embeddings_tensor.cpu().detach().numpy()
        return embeddings_np.astype("float32").tolist()
    
    @traceable(name="query_embedding")
    def embed_query(self, text):
        if self.model is None:
            return [0.0] * 768  # Fallback embedding
        embeddings_tensor = self.model.encode([text], convert_to_tensor=True)
        embeddings_np = embeddings_tensor.cpu().detach().numpy()
        return embeddings_np.astype("float32").tolist()[0]

# Build FAISS index
retriever = None
if doc_chunks and embedding_model is not None:
    try:
        embedding_function = EmbeddingFunction(embedding_model)
        texts = [doc.page_content for doc in doc_chunks]
        embeddings = embedding_model.encode(texts, convert_to_tensor=True).cpu().detach().numpy().astype("float32")
        
        dimension = embeddings.shape[1]
        index = faiss.IndexFlatL2(dimension)
        index.add(embeddings)
        
        docstore = InMemoryDocstore({str(i): doc for i, doc in enumerate(doc_chunks)})
        index_to_docstore_id = {i: str(i) for i in range(len(doc_chunks))}
        
        vectorstore = FAISS(
            embedding_function=embedding_function,
            index=index,
            docstore=docstore,
            index_to_docstore_id=index_to_docstore_id,
        )
        
        retriever = vectorstore.as_retriever()
        print(f"✅ Indexed {len(documents)} files into {len(doc_chunks)} chunks.")
    except Exception as e:
        print(f"❌ Error building FAISS index: {e}")
        retriever = None
else:
    print("⚠️ No documents loaded or embedding model failed, retriever will use fallback")

# ==============================================================================
# LANGUAGE MODEL SETUP AND GRADERS
# ==============================================================================

# Initialize Cohere LLM
try:
    base_llm = ChatCohere(
        model="command-r", 
        temperature=0,
        cohere_api_key=os.environ["COHERE_API_KEY"],
        callbacks=[langchain_tracer] if 'langchain_tracer' in locals() else []
    )
    print("✅ Cohere LLM initialized successfully")
except Exception as e:
    print(f"❌ Error initializing Cohere LLM: {e}")
    base_llm = None

# ==============================================================================
# ROUTER SETUP
# ==============================================================================

class WebSearch(BaseModel):
    """Use for questions NOT related to AI offers, AI foundations, Devoteam practice data, or cybersecurity topics."""
    query: str = Field(description="The query to use when searching the internet.")

class VectorStore(BaseModel):
    """Use for questions related to AI offer generation, AI foundations, Devoteam practice data, and cybersecurity topics."""
    query: str = Field(description="The query to use when searching the vectorstore.")

# Router setup
router_preamble = """
You are an assistant expert in AI cybersecurity offer generation.
Route the question to the appropriate data source:
- Use VectorStore for cybersecurity AI offers, AI foundations, Devoteam practice data
- Use WebSearch for general questions not related to these topics
"""

if base_llm:
    structured_llm_router = base_llm.bind_tools(
        tools=[WebSearch, VectorStore],
        preamble=router_preamble
    )
    
    route_prompt = ChatPromptTemplate.from_messages([("human", "{question}")])
    question_router = route_prompt | structured_llm_router
else:
    question_router = None

# ==============================================================================
# GRADERS SETUP
# ==============================================================================

# Document Relevance Grader
class GradeDocuments(BaseModel):
    binary_score: str = Field(
        description="Documents are relevant to the question, 'yes' or 'no'",
        enum=["yes", "no"]
    )

grade_prompt = ChatPromptTemplate.from_messages([
    ("system", """You are a grader assessing relevance of a retrieved document to a user question.
    
    Your task is to determine if the document contains information that is relevant to answering the user's question.
    
    - If the document contains keywords, concepts, or semantic meaning related to the user question, respond with 'yes'
    - If the document does not contain relevant information, respond with 'no'
    
    You must respond with ONLY 'yes' or 'no' - no other text or explanation."""),
    ("human", "Document: {document}\n\nQuestion: {question}\n\nIs this document relevant to the question? Answer with 'yes' or 'no':")
])

if base_llm:
    structured_llm_grader = base_llm.with_structured_output(GradeDocuments)
    retrieval_grader = grade_prompt | structured_llm_grader
else:
    retrieval_grader = None

# Hallucination Grader
class GradeHallucinations(BaseModel):
    """Binary score for hallucination present in generation answer."""
    binary_score: str = Field(
        description="Answer is grounded in the facts, 'yes' or 'no'",
        enum=["yes", "no"]
    )

hallucination_prompt = ChatPromptTemplate.from_messages([
    ("system", """You are a grader assessing whether an LLM generation is grounded in / supported by a set of retrieved facts.
    
    Your task is to determine if the generation is factually grounded in the provided documents.
    
    - If the generation is supported by the documents, respond with 'yes'
    - If the generation contains information not found in the documents, respond with 'no'
    
    You must respond with ONLY 'yes' or 'no' - no other text or explanation."""),
    ("human", "Documents: {documents}\n\nLLM Generation: {generation}\n\nIs the generation grounded in the documents? Answer with 'yes' or 'no':")
])

if base_llm:
    structured_hallucination_grader = base_llm.with_structured_output(GradeHallucinations)
    hallucination_grader = hallucination_prompt | structured_hallucination_grader
else:
    hallucination_grader = None

# Answer Quality Grader
class GradeAnswer(BaseModel):
    """Binary score to assess answer addresses question."""
    binary_score: str = Field(
        description="Answer addresses the question, 'yes' or 'no'",
        enum=["yes", "no"]
    )

answer_prompt = ChatPromptTemplate.from_messages([
    ("system", """You are a grader assessing whether an answer addresses / resolves a question.
    Your task is to determine if the answer properly addresses the user's question.
    
    - If the answer addresses and resolves the question, respond with 'yes'
    - If the answer does not address or resolve the question, respond with 'no'
    
    You must respond with ONLY 'yes' or 'no' - no other text or explanation."""),
    ("human", "Question: {question}\n\nAnswer: {generation}\n\nDoes the answer address the question? Answer with 'yes' or 'no':")
])

if base_llm:
    structured_answer_grader = base_llm.with_structured_output(GradeAnswer)
    answer_grader = answer_prompt | structured_answer_grader
else:
    answer_grader = None

# ==============================================================================
# GENERATION FUNCTIONS
# ==============================================================================

# Main generation function
generation_preamble = """You are an assistant for cybersecurity AI offer generation. Use the following pieces of retrieved context to generate a comprehensive cybersecurity offer that includes:
1. RÉSUMÉ EXÉCUTIF: Brief overview of the AI-powered cybersecurity solution
2. SOLUTION TECHNIQUE: Detailed approach using AI for cybersecurity
3. PLAN DE MISE EN ŒUVRE: Step-by-step deployment strategy
4. CALENDRIER: Project phases and milestones
5. COMPOSITION DE L'ÉQUIPE: Required expertise and roles
6. RÉSULTATS ATTENDUS: Measurable benefits and ROI
7. STRUCTURE DE PRIX: Cost breakdown and investment

Generate a professional, detailed offer in French. If you don't know specific details, use the context provided and general cybersecurity AI best practices."""

if base_llm:
    generation_llm = base_llm.bind(preamble=generation_preamble)
else:
    generation_llm = None

def enhanced_prompt(x):
    return ChatPromptTemplate.from_messages([
        HumanMessage(
            f"Question: {x['question']}\nContext: {x['documents']}\n\nGenerate a comprehensive cybersecurity AI offer:",
            additional_kwargs={"documents": x["documents"]},
        )
    ])

@traceable(name="generate_basic")
def generate_basic(state):
    """Generate answer with documents"""
    print("---GENERATE---")
    question = state["question"]
    documents = state["documents"]
    
    if generation_llm is None:
        return {
            "question": question,
            "documents": documents,
            "generation": "Erreur: Le modèle de langage n'est pas disponible. Veuillez vérifier votre configuration Cohere."
        }
    
    try:
        # Convert documents to text
        if isinstance(documents, list) and documents:
            documents_text = "\n\n".join([
                doc.page_content if hasattr(doc, 'page_content') else str(doc) 
                for doc in documents
            ])
        else:
            documents_text = str(documents) if documents else ""
        
        # Generate the offer
        rag_chain = enhanced_prompt | generation_llm | StrOutputParser()
        generation = rag_chain.invoke({"documents": documents_text, "question": question})
        
        return {
            "question": question,
            "documents": documents,
            "generation": generation
        }
    except Exception as e:
        print(f"Error in generate_basic: {e}")
        return {
            "question": question,
            "documents": documents,
            "generation": f"Erreur lors de la génération: {str(e)}"
        }

# LLM Fallback function
fallback_preamble = """You are an assistant for question-answering tasks. Answer the question based upon your knowledge. Use three sentences maximum and keep the answer concise."""

if base_llm:
    fallback_llm = base_llm.bind(preamble=fallback_preamble)
    
    def fallback_prompt(x):
        return ChatPromptTemplate.from_messages([
            HumanMessage(f"Question: {x['question']} \nAnswer: ")
        ])
    
    llm_chain = fallback_prompt | fallback_llm | StrOutputParser()
else:
    llm_chain = None

# Dummy web search function (since web search is not available in your setup)
def dummy_web_search(query):
    """Dummy web search function"""
    return [Document(
        page_content=f"Web search results for: {query}. This is a placeholder response.",
        metadata={"source": "web_search"}
    )]

# run the question directly on the LLM
def generate_llm_only_answer(question):
    """Generate an answer using the LLM without RAG context"""
    if llm_chain is None:
        return "Le modèle LLM n'est pas disponible."
    try:
        return llm_chain.invoke({"question": question})
    except Exception as e:
        return f"Erreur LLM direct: {str(e)}"



# similarity calculation between LLM-only and RAG outputs
def calculate_similarity(text1, text2):
    """Calculate cosine similarity between two texts using embeddings"""
    if not text1 or not text2 or embedding_model is None:
        return 0.0
    try:
        emb1 = embedding_model.encode([text1], convert_to_tensor=True).cpu().detach().numpy()
        emb2 = embedding_model.encode([text2], convert_to_tensor=True).cpu().detach().numpy()
        return float(cosine_similarity(emb1, emb2)[0][0])
    except Exception as e:
        print(f"Error calculating similarity: {e}")
        return 0.0

# ==============================================================================
# WORKFLOW FUNCTIONS
# ==============================================================================

class GraphState(TypedDict):
    question: str
    generation: str
    documents: List[Document]

def retrieve(state):
    """Retrieve documents"""
    print("---RETRIEVE---")
    question = state["question"]
    
    if retriever is None:
        print("No retriever available, using empty documents")
        return {"documents": [], "question": question}
    
    try:
        documents = retriever.invoke(question)
        print(f"Retrieved {len(documents)} documents")
        return {"documents": documents, "question": question}
    except Exception as e:
        print(f"Error during retrieval: {e}")
        return {"documents": [], "question": question}

def llm_fallback(state):
    """Generate answer using the LLM w/o vectorstore"""
    print("---LLM Fallback---")
    question = state["question"]
    
    if llm_chain is None:
        generation = "Le service de génération n'est pas disponible. Veuillez vérifier la configuration."
    else:
        try:
            generation = llm_chain.invoke({"question": question})
        except Exception as e:
            generation = f"Erreur lors de la génération: {str(e)}"
    
    return {"question": question, "generation": generation}

def generate(state):
    """Generate answer using the vectorstore"""
    return generate_basic(state)

def grade_documents(state):
    """Determines whether the retrieved documents are relevant to the question."""
    print("---CHECK DOCUMENT RELEVANCE TO QUESTION---")
    question = state["question"]
    documents = state["documents"]
    
    if not documents or retrieval_grader is None:
        print("---NO DOCUMENTS OR GRADER AVAILABLE---")
        return {"documents": documents, "question": question}

    filtered_docs = []
    try:
        for d in documents:
            score = retrieval_grader.invoke(
                {"question": question, "document": d.page_content}
            )
            grade = score.binary_score
            if grade == "yes":
                print("---GRADE: DOCUMENT RELEVANT---")
                filtered_docs.append(d)
            else:
                print("---GRADE: DOCUMENT NOT RELEVANT---")
                continue
    except Exception as e:
        print(f"Error in document grading: {e}")
        filtered_docs = documents  # Use all documents if grading fails
    
    return {"documents": filtered_docs, "question": question}

def web_search(state):
    """Web search based on the question."""
    print("---WEB SEARCH---")
    question = state["question"]
    
    # Using dummy web search since actual web search is not available
    docs = dummy_web_search(question)
    
    return {"documents": docs, "question": question}

def route_question(state):
    """Route question to web search or RAG."""
    print("---ROUTE QUESTION---")
    question = state["question"]
    
    if question_router is None:
        print("---NO ROUTER AVAILABLE, DEFAULTING TO VECTORSTORE---")
        return "vectorstore"
    
    try:
        source = question_router.invoke({"question": question})

        if hasattr(source, 'tool_calls') and source.tool_calls:
            tool_call = source.tool_calls[0]
            datasource = tool_call.get('name', 'VectorStore')
        elif hasattr(source, 'response_metadata') and source.response_metadata.get('tool_calls'):
            tool_calls = source.response_metadata['tool_calls']
            if tool_calls:
                datasource = tool_calls[0].get('function', {}).get('name', 'VectorStore')
            else:
                datasource = 'VectorStore'
        else:
            print("---ROUTE QUESTION TO LLM---")
            return "llm_fallback"

        if datasource == "WebSearch":
            print("---ROUTE QUESTION TO WEB SEARCH---")
            return "web_search"
        elif datasource == "VectorStore":
            print("---ROUTE QUESTION TO RAG---")
            return "vectorstore"
        else:
            print("---ROUTE QUESTION TO RAG---")
            return "vectorstore"
    except Exception as e:
        print(f"Error in routing: {e}, defaulting to vectorstore")
        return "vectorstore"

def decide_to_generate(state):
    """Determines whether to generate an answer, or do web search."""
    print("---ASSESS GRADED DOCUMENTS---")
    filtered_documents = state["documents"]

    if not filtered_documents:
        print("---DECISION: ALL DOCUMENTS ARE NOT RELEVANT TO QUESTION, WEB SEARCH---")
        return "web_search"
    else:
        print("---DECISION: GENERATE---")
        return "generate"

def grade_generation_v_documents_and_question(state):
    """Determines whether the generation is grounded in the document and answers question."""
    print("---CHECK HALLUCINATIONS---")
    question = state["question"]
    documents = state["documents"]
    generation = state["generation"]

    if hallucination_grader is None or answer_grader is None:
        print("---GRADERS NOT AVAILABLE, ASSUMING USEFUL---")
        return "useful"

    try:
        # Convert documents to text for grading
        if isinstance(documents, list) and documents:
            documents_text = "\n\n".join([
                doc.page_content if hasattr(doc, 'page_content') else str(doc) 
                for doc in documents
            ])
        else:
            documents_text = str(documents) if documents else ""

        score = hallucination_grader.invoke(
            {"documents": documents_text, "generation": generation}
        )
        grade = score.binary_score

        if grade == "yes":
            print("---DECISION: GENERATION IS GROUNDED IN DOCUMENTS---")
            print("---GRADE GENERATION vs QUESTION---")
            score = answer_grader.invoke({"question": question, "generation": generation})
            grade = score.binary_score
            if grade == "yes":
                print("---DECISION: GENERATION ADDRESSES QUESTION---")
                return "useful"
            else:
                print("---DECISION: GENERATION DOES NOT ADDRESS QUESTION---")
                return "not useful"
        else:
            print("---DECISION: GENERATION IS NOT GROUNDED IN DOCUMENTS, RE-TRY---")
            return "not supported"
    except Exception as e:
        print(f"Error in generation grading: {e}, assuming useful")
        return "useful"

# ==============================================================================
# METRICS CALCULATION FUNCTIONS
# ==============================================================================

def calculate_quality_metrics(generation, documents, question):
    """Calculate comprehensive quality metrics for the generated offer"""
    metrics = {
        "quality": {
            "overall_score": 0.0,  # Will be calculated
            "grade": "F",  # Default to F, will be updated
            "detailed_scores": {
                "precision": 0.0,
                "consistency": 0.0,
                "completeness": 0.0,
                "relevance": 0.0,
            }
        }
    }
    
    try:
        # 1. Calculate Precision (technical term coverage)
        technical_terms = [
            "cybersécurité", "sécurité", "vulnérabilité", "menace",
            "détection", "protection", "conformité", "audit", "risque",
            "IA", "intelligence artificielle", "machine learning"
        ]
        if generation:
            found_terms = sum(
                1 for term in technical_terms 
                if term.lower() in generation.lower()
            )
            metrics["quality"]["detailed_scores"]["precision"] = min(
                found_terms / len(technical_terms), 
                1.0
            )

        # 2. Calculate Consistency (structural quality)
        if generation:
            sentences = [s.strip() for s in re.split(r'[.!?]', generation) if s.strip()]
            if sentences:
                avg_length = sum(len(s.split()) for s in sentences) / len(sentences)
                # Ideal range 15-25 words per sentence
                metrics["quality"]["detailed_scores"]["consistency"] = min(
                    max((avg_length - 10) / 15, 0),  # Normalize 10-25 words to 0-1
                    1.0
                )

        # 3. Calculate Completeness (required sections)
        section_groups = [
            ["RÉSUMÉ EXÉCUTIF", "EXECUTIVE SUMMARY"],
            ["SOLUTION TECHNIQUE", "TECHNICAL SOLUTION"],
            ["PLAN DE MISE EN ŒUVRE", "IMPLEMENTATION PLAN"],
            ["CALENDRIER", "TIMELINE", "SCHEDULE"],
            ["COMPOSITION DE L'ÉQUIPE", "TEAM COMPOSITION"],
            ["RÉSULTATS ATTENDUS", "EXPECTED RESULTS"],
            ["STRUCTURE DE PRIX", "PRICING STRUCTURE"]
        ]
        
        if generation:
            gen_lower = generation.lower()
            found_sections = 0
            
            for group in section_groups:
                if any(section.lower() in gen_lower for section in group):
                    found_sections += 1
            
            metrics["quality"]["detailed_scores"]["completeness"] = found_sections / len(section_groups)

        # 4. Calculate Relevance (question alignment)
        if question and generation:
            question_words = set(word.lower() for word in question.split() if len(word) > 3)
            generation_words = set(word.lower() for word in generation.split() if len(word) > 3)
            common_words = question_words.intersection(generation_words)
            
            if question_words:
                metrics["quality"]["detailed_scores"]["relevance"] = len(common_words) / len(question_words)

        # Calculate weighted overall score
        weights = {
            "precision": 0.3,
            "consistency": 0.2,
            "completeness": 0.3,
            "relevance": 0.2
        }
        
        total_score = sum(
            weight * metrics["quality"]["detailed_scores"][metric]
            for metric, weight in weights.items()
        )
        metrics["quality"]["overall_score"] = min(max(total_score, 0), 1)  # Clamp to 0-1

        # Determine letter grade
        overall = metrics["quality"]["overall_score"]
        if overall >= 0.9:
            metrics["quality"]["grade"] = "A"
        elif overall >= 0.8:
            metrics["quality"]["grade"] = "B"
        elif overall >= 0.7:
            metrics["quality"]["grade"] = "C"
        elif overall >= 0.6:
            metrics["quality"]["grade"] = "D"

    except Exception as e:
        print(f"Error calculating quality metrics: {e}")
        # Fallback to default metrics if error occurs
    
    return metrics

def calculate_cost_metrics(generation, question):
    """Calculate cost and token usage metrics"""
    metrics = {
        "cost": 0.0,
        "input_tokens": len(question.split()) if question else 0,
        "output_tokens": len(generation.split()) if generation else 0
    }
    
    try:
        # Cohere pricing estimate (adjust based on your model)
        input_cost_per_token = 0.0015 / 1000  # $0.0015 per 1000 input tokens
        output_cost_per_token = 0.0020 / 1000  # $0.0020 per 1000 output tokens
        
        metrics["cost"] = round(
            (metrics["input_tokens"] * input_cost_per_token) +
            (metrics["output_tokens"] * output_cost_per_token),
            6
        )
    except Exception as e:
        print(f"Error calculating cost metrics: {e}")
    
    return metrics

def calculate_performance_metrics(start_time, end_time, documents_retrieved, relevant_documents):
    """Calculate performance metrics"""
    return {
        "processing_time": round(end_time - start_time, 2),
        "memory_usage": psutil.Process().memory_info().rss / (1024 * 1024),  # MB
        "documents_retrieved": documents_retrieved,
        "relevant_documents": relevant_documents,
        "relevance_ratio": (
            (relevant_documents / documents_retrieved * 100) 
            if documents_retrieved > 0 else 0
        )
    }
# ==============================================================================
# WORKFLOW GRAPH
# ==============================================================================

try:
    workflow = StateGraph(GraphState)

    # Define the nodes
    workflow.add_node("web_search", web_search)
    workflow.add_node("retrieve", retrieve)
    workflow.add_node("grade_documents", grade_documents)
    workflow.add_node("generate", generate)
    workflow.add_node("llm_fallback", llm_fallback)

    # Build graph with conditional edges
    workflow.add_conditional_edges(
        START,
        route_question,
        {
            "web_search": "web_search",
            "vectorstore": "retrieve",
            "llm_fallback": "llm_fallback",
        },
    )

    workflow.add_edge("web_search", "generate")
    workflow.add_edge("retrieve", "grade_documents")

    workflow.add_conditional_edges(
        "grade_documents",
        decide_to_generate,
        {
            "web_search": "web_search",
            "generate": "generate",
        },
    )

    workflow.add_conditional_edges(
        "generate",
        grade_generation_v_documents_and_question,
        {
            "not supported": "generate",
            "not useful": "web_search",
            "useful": END,
        },
    )

    workflow.add_edge("llm_fallback", END)

    # Compile the workflow
    app_graph = workflow.compile()
    print("✅ Enhanced workflow compiled successfully.")
except Exception as e:
    print(f"❌ Error compiling enhanced workflow: {e}")
    app_graph = None

# ==============================================================================
# FLASK ROUTES
# ==============================================================================

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate_offer():
    try:
        data = request.get_json()
        question = data.get('question', '').strip()
        
        if not question:
            return jsonify({'error': 'Question is required'}), 400
        
        print(f"🔍 Processing query: {question}")
        
        result = {
            "steps": [], 
            "final_result": None, 
            "metrics": {},
            "documents_retrieved": 0,
            "relevant_documents": 0,
            "grading_details": {
                "document_relevance": [],
                "hallucination_check": None,
                "answer_quality": None
            }
        }
        
        start_time = time.time()
        
        if app_graph is None:
            # Fallback without workflow
            print("⚠️ Using fallback generation")
            fallback_result = llm_fallback({"question": question})
            result["final_result"] = fallback_result["generation"]
        else:
            # Execute the enhanced workflow
            inputs = {"question": question}
            final_state = None
            
            for output in app_graph.stream(inputs):
                for key, value in output.items():
                    print(f"🔄 Pipeline step: {key}")
                    
                    step_info = {
                        "step": key,
                        "status": "completed",
                        "timestamp": time.time(),
                        "description": get_step_description(key)
                    }
                    
                    # Track document counts
                    if "documents" in value and value["documents"]:
                        documents_count = len(value["documents"])
                        step_info["documents_count"] = documents_count
                        
                        if key == "retrieve":
                            result["documents_retrieved"] = documents_count
                        elif key == "grade_documents":
                            result["relevant_documents"] = documents_count
                            # Track document grading details
                            for doc in value["documents"]:
                                result["grading_details"]["document_relevance"].append({
                                    "content_preview": doc.page_content[:100] + "..." if hasattr(doc, 'page_content') else str(doc)[:100] + "...",
                                    "relevant": True
                                })
                    
                    # Track generation
                    if "generation" in value and value["generation"]:
                        result["final_result"] = value["generation"]
                    
                

                    result["steps"].append(step_info)
                    final_state = value
        
        end_time = time.time()

        if result["final_result"]:
            # Run direct LLM (no RAG)
            llm_direct_answer = generate_llm_only_answer(question)

            # Calculate similarity between RAG and direct LLM answers
            similarity_score = calculate_similarity(result["final_result"], llm_direct_answer)

            # Save in result
            result["comparison"] = {
                            "llm_only_answer": llm_direct_answer,
                            "rag_answer": result["final_result"],
                            "similarity_score": similarity_score
                        }

        
        # Ensure we have a result
        if not result["final_result"]:
            result["final_result"] = "Erreur: Aucune génération d'offre n'a été produite."
        
        # Calculate all metrics
        performance_metrics = calculate_performance_metrics(
            start_time, end_time, 
            result["documents_retrieved"], 
            result["relevant_documents"]
        )
        
        quality_metrics = calculate_quality_metrics(
            result["final_result"],
            final_state["documents"] if final_state and "documents" in final_state else [],
            question
        )
        
        cost_metrics = calculate_cost_metrics(result["final_result"], question)
        
        # Combine all metrics
        result["metrics"] = {
            **performance_metrics,
            **quality_metrics,
            **cost_metrics,
            "taxonomy": classify_offer_taxonomy(question),
            "confidence_score": calculate_confidence_score(result),
            "pipeline_efficiency": calculate_pipeline_efficiency(result["steps"]),
            "workflow_path": [step["step"] for step in result["steps"]]
        }
        
        print(f"✅ Enhanced pipeline completed in {result['metrics']['processing_time']}s")
        
        return jsonify(result)
        
    except Exception as e:
        error_msg = f"Error in generate_offer: {str(e)}"
        print(f"❌ {error_msg}")
        print(traceback.format_exc())
        return jsonify({'error': error_msg}), 500

def get_step_description(step_name):
    """Get human-readable description for pipeline steps"""
    descriptions = {
        "retrieve": "Récupération des documents pertinents depuis la base de connaissances",
        "grade_documents": "Évaluation de la pertinence des documents récupérés",
        "generate": "Génération de l'offre cybersécurité IA personnalisée",
        "web_search": "Recherche d'informations complémentaires sur le web",
        "llm_fallback": "Génération de réponse avec le modèle de langage de base"
    }
    return descriptions.get(step_name, f"Exécution de l'étape: {step_name}")

def calculate_confidence_score(result):
    """Calculate confidence score based on pipeline execution"""
    score = 50  # Base score
    
    # Add points for document retrieval
    if result["documents_retrieved"] > 0:
        score += 15
    if result["relevant_documents"] > 0:
        score += 20
    
    # Add points for pipeline completion
    if result["final_result"] and len(result["final_result"]) > 100:
        score += 10
    
    # Add points for workflow efficiency
    if len(result["steps"]) <= 4:  # Efficient workflow
        score += 5
    
    return min(score, 95)

def calculate_pipeline_efficiency(steps):
    """Calculate pipeline efficiency based on number of steps and retries"""
    if not steps:
        return 0
    
    total_steps = len(steps)
    unique_steps = len(set(step["step"] for step in steps))
    
    # Calculate efficiency: fewer total steps relative to unique steps is better
    if unique_steps == 0:
        return 0
    
    efficiency = (unique_steps / total_steps) * 100
    return round(efficiency, 2)

@app.route('/health')
def health_check():
    """Enhanced health check with component status"""
    component_status = {
        'status': 'healthy',
        'components': {
            'documents_loaded': len(documents),
            'chunks_available': len(doc_chunks) if doc_chunks else 0,
            'retriever_ready': retriever is not None,
            'langsmith_configured': bool(os.environ.get("LANGCHAIN_API_KEY")),
            'llm_ready': base_llm is not None,
            'embedding_model_ready': embedding_model is not None,
            'router_ready': question_router is not None,
            'graders_ready': {
                'retrieval_grader': retrieval_grader is not None,
                'hallucination_grader': hallucination_grader is not None,
                'answer_grader': answer_grader is not None
            },
            'workflow_ready': app_graph is not None
        },
        'taxonomy_categories': list(OFFER_TAXONOMY.keys()),
        'system_capabilities': {
            'document_processing': True,
            'semantic_search': retriever is not None,
            'ai_generation': base_llm is not None,
            'quality_control': all([
                retrieval_grader is not None,
                hallucination_grader is not None,
                answer_grader is not None
            ]),
            'workflow_routing': question_router is not None
        }
    }
    
    # Determine overall health
    critical_components = [
        component_status['components']['llm_ready'],
        component_status['components']['workflow_ready']
    ]
    
    if not all(critical_components):
        component_status['status'] = 'degraded'
    
    return jsonify(component_status)

@app.route('/test-components', methods=['GET'])
def test_components():
    """Test individual components"""
    test_results = {}
    test_question = "Test de génération d'offre IA cybersécurité"
    
    # Test retriever
    if retriever:
        try:
            docs = retriever.invoke(test_question)
            test_results['retriever'] = {
                'status': 'working',
                'documents_found': len(docs),
                'first_doc_preview': docs[0].page_content[:100] + "..." if docs else "No content"
            }
        except Exception as e:
            test_results['retriever'] = {'status': 'error', 'error': str(e)}
    else:
        test_results['retriever'] = {'status': 'not_available'}
    
    # Test router
    if question_router:
        try:
            route_result = question_router.invoke({"question": test_question})
            test_results['router'] = {
                'status': 'working',
                'result': str(route_result)[:200] + "..."
            }
        except Exception as e:
            test_results['router'] = {'status': 'error', 'error': str(e)}
    else:
        test_results['router'] = {'status': 'not_available'}
    
    # Test retrieval grader
    if retrieval_grader and retriever:
        try:
            docs = retriever.invoke(test_question)
            if docs:
                grade_result = retrieval_grader.invoke({
                    "question": test_question, 
                    "document": docs[0].page_content
                })
                test_results['retrieval_grader'] = {
                    'status': 'working',
                    'grade': grade_result.binary_score
                }
            else:
                test_results['retrieval_grader'] = {'status': 'no_documents_to_grade'}
        except Exception as e:
            test_results['retrieval_grader'] = {'status': 'error', 'error': str(e)}
    else:
        test_results['retrieval_grader'] = {'status': 'not_available'}
    
    # Test generation
    if base_llm:
        try:
            test_state = {
                "question": test_question,
                "documents": ["Test document content for cybersecurity AI offer generation."]
            }
            gen_result = generate_basic(test_state)
            test_results['generation'] = {
                'status': 'working',
                'output_length': len(gen_result.get('generation', '')),
                'preview': gen_result.get('generation', '')[:150] + "..."
            }
        except Exception as e:
            test_results['generation'] = {'status': 'error', 'error': str(e)}
    else:
        test_results['generation'] = {'status': 'not_available'}
    
    # Test workflow
    if app_graph:
        try:
            # Run a simple workflow test
            inputs = {"question": "Test simple"}
            workflow_steps = []
            for output in app_graph.stream(inputs):
                workflow_steps.extend(output.keys())
                break  # Just test first step
            
            test_results['workflow'] = {
                'status': 'working',
                'first_step': workflow_steps[0] if workflow_steps else 'no_steps'
            }
        except Exception as e:
            test_results['workflow'] = {'status': 'error', 'error': str(e)}
    else:
        test_results['workflow'] = {'status': 'not_available'}
    
    return jsonify({
        'test_question': test_question,
        'timestamp': time.time(),
        'results': test_results
    })

@app.route('/taxonomy', methods=['GET'])
def get_taxonomy():
    """Get the offer taxonomy structure"""
    return jsonify({
        'taxonomy': OFFER_TAXONOMY,
        'categories': list(OFFER_TAXONOMY.keys()),
        'total_subcategories': sum(len(subcats) for subcats in OFFER_TAXONOMY.values())
    })

@app.route('/taxonomy/classify', methods=['POST'])
def classify_question():
    """Classify a question using the taxonomy"""
    try:
        data = request.get_json()
        question = data.get('question', '').strip()
        
        if not question:
            return jsonify({'error': 'Question is required'}), 400
        
        classification = classify_offer_taxonomy(question)
        return jsonify({
            'question': question,
            'classification': classification,
            'timestamp': time.time()
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/documents/stats', methods=['GET'])
def get_document_stats():
    """Get document processing statistics"""
    stats = {
        'total_documents': len(documents),
        'total_chunks': len(doc_chunks) if doc_chunks else 0,
        'average_chunk_size': 0,
        'document_sources': [],
        'processing_status': {
            'documents_loaded': len(documents) > 0,
            'chunks_created': len(doc_chunks) > 0 if doc_chunks else False,
            'embeddings_created': retriever is not None,
            'vectorstore_ready': retriever is not None
        }
    }
    
    if doc_chunks:
        total_length = sum(len(chunk.page_content) for chunk in doc_chunks)
        stats['average_chunk_size'] = round(total_length / len(doc_chunks), 2)
    
    # Get document sources (filenames only for privacy)
    for doc in documents:
        if 'source' in doc.metadata:
            filename = os.path.basename(doc.metadata['source'])
            stats['document_sources'].append(filename)
    
    return jsonify(stats)

# ==============================================================================
# ERROR HANDLERS
# ==============================================================================

@app.errorhandler(404)
def not_found(error):
    return jsonify({'error': 'Endpoint not found'}), 404

@app.errorhandler(500)
def internal_error(error):
    return jsonify({'error': 'Internal server error'}), 500

@app.errorhandler(Exception)
def handle_exception(e):
    # Log the error for debugging
    print(f"Unhandled exception: {str(e)}")
    print(traceback.format_exc())
    return jsonify({'error': 'An unexpected error occurred'}), 500

# ==============================================================================
# MAIN APPLICATION STARTUP
# ==============================================================================

if __name__ == '__main__':
    print("\n" + "="*80)
    print("🚀 Starting Enhanced Cybersecurity AI Offer Generator...")
    print("="*80)
    
    # System status summary
    print("📊 SYSTEM STATUS:")
    print(f"   📄 Documents loaded: {len(documents)}")
    print(f"   📚 Chunks available: {len(doc_chunks) if doc_chunks else 0}")
    print(f"   🔍 Retriever ready: {retriever is not None}")
    print(f"   🤖 LLM ready: {base_llm is not None}")
    print(f"   🧠 Embedding model ready: {embedding_model is not None}")
    print(f"   🎯 Router ready: {question_router is not None}")
    print(f"   ✅ Workflow ready: {app_graph is not None}")
    
    print("\n📋 COMPONENT STATUS:")
    print(f"   📝 Retrieval Grader: {retrieval_grader is not None}")
    print(f"   🔍 Hallucination Grader: {hallucination_grader is not None}")
    print(f"   ⭐ Answer Grader: {answer_grader is not None}")
    print(f"   🌐 LangSmith Tracing: {'langchain_tracer' in locals()}")
    
    print("\n🏗️ AVAILABLE ENDPOINTS:")
    print("   GET  /              - Main interface")
    print("   POST /generate      - Generate cybersecurity offers")
    print("   GET  /health        - System health check")
    print("   GET  /test-components - Test individual components")
    print("   GET  /taxonomy      - Get taxonomy structure")
    print("   POST /taxonomy/classify - Classify questions")
    print("   GET  /documents/stats - Document statistics")
    
    print("="*80)
    print("🌐 Server starting on http://localhost:5000")
    print("📖 Visit /health for system status")
    print("🧪 Visit /test-components for component testing")
    print("="*80)
    
    try:
        app.run(debug=True, host='0.0.0.0', port=5000, threaded=True)
    except Exception as e:
        print(f"❌ Failed to start Flask server: {e}")
        print("Make sure port 5000 is not in use by another application.")
        print("You can also try running on a different port by changing the port parameter.")